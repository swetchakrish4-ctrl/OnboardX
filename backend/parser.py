

from pypdf import PdfReader
from pptx import Presentation
import email

def extract_text(file):
    name = file.name.lower()
    if name.endswith('.pdf'):
        return " ".join([p.extract_text() for p in PdfReader(file).pages])
    elif name.endswith('.pptx'):
        text = ""
        for slide in Presentation(file).slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"): text += shape.text + " "
        return text
    
    return ""